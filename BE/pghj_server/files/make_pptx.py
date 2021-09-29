from pptx import Presentation
from pptx.util import Inches, Pt

from pghj_server.settings import PPTX_DIR # ppt 저장경로


# 기존에 존재하는 템플릿 수정해서 pptx 만들기
def update_pptx(prs, items):
    for slide in prs.slides:
        # print(slide.slide_id) # 256, 257, 258
        idx = slide.slide_id - prs.slides[0].slide_id + 1
        
        shapes = slide.shapes
        # print(idx)

        while items != []:
            item = items[0]

            page = item['page']
            text = item['text']
            left = Inches(item['coordinate']['left'])
            top = Inches(item['coordinate']['top'])
            width = Inches(item['size']['width'])
            height = Inches(item['size']['height'])
            font_type = item['font']['type']
            font_size = item['font']['size']            

            if page > idx: # 다음 페이지의 item이 등장하면 루프 멈추고 다음 슬라이드 불러오기
                break
            
            text_box = shapes.add_textbox(left, top, width, height) # 텍스트박스 만들기
            text_frame = text_box.text_frame # 텍스트 박스의 프레임 불러오기
            p = text_frame.add_paragraph()
            p.text = text

            if font_type == "bold": p.font.bold = True
            if font_type == "italic": p.font.italic = True
            p.font_size = Pt(font_size)      

            items.pop(0)
    return prs



def get_pptx(data, user_id, upload_id):  
    template_id = data['template_id'] # 템플릿 디자인
    items = data['items']             # 유저 선택 옵션 dict_type으로 받기        

    if template_id == 0:              # 만약 template id가 0이면 아예 빈 슬라이드 만들기(페이지 수만큼)
        prs = Presentation()
        last_page = items[len(items)-1]['page']
        for _ in range(last_page):
            slide_layout = prs.slide_layouts[6] # 레이아웃 불러오기
            prs.slides.add_slide(slide_layout)  # 해당 레이아웃을 가진 슬라이드 추가하기

    else:                             # template_id가 있으면 기존 템플릿 불러오기(페이지 수 조절 불가. 그냥 템플릿 위에 씌우기만 가능)
        prs = Presentation(PPTX_DIR + "ppt_templates/" + str(template_id) + ".pptx")

    prs = update_pptx(prs, items) # PPT 만들기

    pptx_path = PPTX_DIR + 'pptx_results/' + user_id + '_' + str(upload_id) + '.pptx'

    prs.save(pptx_path)   # ppt 저장하기

    return pptx_path
