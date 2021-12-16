from pptx import Presentation
from pptx.util import Inches, Pt
from pghj_server.settings import TEMPLATE_DIR

# Change pixel into inch
def get_inch(x, y):
    r = x * 1280 / 96
    c = y * 720 / 96
    return Inches(r), Inches(c)
  
# Add slides
def add_slides(prs, items):
    # add each slide
    for item in items:                             # one item per page
        slide_layout = prs.slide_layouts[6]        # create slide_layout (6: BLANK)
        slide = prs.slides.add_slide(slide_layout) # add a slide
        shapes = slide.shapes

        sentences = item['sentences']
        for sentence_block in sentences:
            sentence = sentence_block['sentence']
            left, top = get_inch(sentence_block['coordinate']['left'], sentence_block['coordinate']['top'])
            width, height = get_inch(sentence_block['size']['width'], sentence_block['size']['height'])
            font_type = sentence_block['font']['type']
            font_size = sentence_block['font']['size']

            text_box = shapes.add_textbox(left, top, width, height) # Create text box
            text_frame = text_box.text_frame                        # Get text frame from the text box

            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = sentence                                     # Write text

            font = run.font
            font.name = font_type                                   # Change font type
            font.size = Pt(font_size)                               # Change font size
    return prs


# make presentation
def create_pptx(template_id, file_path, file_name, items):  
    # open presentation
    if template_id == "template00": # If user do not choose any template, open new presentation
        prs = Presentation()
    else:                           # Else, open the presentation with certain template
        prs = Presentation(TEMPLATE_DIR + str(template_id) + '.pptx')
    
    # add slides
    prs = add_slides(prs, items)

    # save pptx
    prs.save(file_path + file_name)